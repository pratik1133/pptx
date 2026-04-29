from pathlib import Path
from uuid import uuid4

from pptx import Presentation

from reportgen.ai.planner import plan_slides_mock
from reportgen.ingestion.loaders import load_normalized_input_bundle
from reportgen.planning.report_builder import build_report_spec
from reportgen.rendering.engine import PresentationRenderer


def test_renderer_writes_pptx() -> None:
    bundle = load_normalized_input_bundle(Path("data/samples/bundles/abc_bundle.json"))
    plan = plan_slides_mock(bundle)
    report_spec = build_report_spec(bundle, plan)

    renderer = PresentationRenderer()
    temp_dir = Path("output/test-temp") / uuid4().hex
    temp_dir.mkdir(parents=True, exist_ok=False)
    out_path = temp_dir / "sample.pptx"

    written = renderer.render_to_path(report_spec, out_path)

    assert written.exists()
    assert written.stat().st_size > 0

    presentation = Presentation(written)
    assert len(presentation.slides) == len(report_spec.slides)

    chart_slide = presentation.slides[3]
    assert any(shape.shape_type == 13 for shape in chart_slide.shapes)

    table_slide = presentation.slides[4]
    assert any(getattr(shape, "has_table", False) for shape in table_slide.shapes)
