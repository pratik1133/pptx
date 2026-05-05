from __future__ import annotations

from pathlib import Path

from pptx import Presentation

from reportgen.qa.validators import QaResult
from reportgen.schemas.charts import ChartBlock
from reportgen.schemas.report import ReportSpec
from reportgen.schemas.tables import TableBlock


def validate_rendered_pptx(pptx_path: Path, report_spec: ReportSpec) -> QaResult:
    result = QaResult()
    if not pptx_path.exists() or pptx_path.stat().st_size == 0:
        result.errors.append("Rendered PPTX artifact is missing or empty.")
        return result

    presentation = Presentation(pptx_path)
    if len(presentation.slides) != len(report_spec.slides):
        result.errors.append(
            f"Rendered PPTX slide count {len(presentation.slides)} does not match report spec {len(report_spec.slides)}."
        )
        return result

    for index, slide_spec in enumerate(report_spec.slides):
        slide = presentation.slides[index]
        if any(isinstance(block, ChartBlock) for block in slide_spec.blocks):
            if not any(shape.shape_type == 13 for shape in slide.shapes):
                result.errors.append(f"Slide {slide_spec.slide_id} is missing a rendered chart image.")
        if any(isinstance(block, TableBlock) for block in slide_spec.blocks):
            if slide_spec.layout not in ("forensic_assessment", "saarthi_scorecard"):
                if not any(getattr(shape, "has_table", False) for shape in slide.shapes):
                    result.errors.append(f"Slide {slide_spec.slide_id} is missing a rendered table.")

    return result
