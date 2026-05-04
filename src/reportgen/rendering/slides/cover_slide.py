from __future__ import annotations

from pathlib import Path
from typing import Any

from reportgen.rendering.brand_shell import (
    _add_filled_rect,
    _add_text_in_rect,
    _format_currency_compact,
    _format_date_ddmmyyyy,
    _format_market_cap,
    _hex_to_rgb,
    apply_background,
)
from reportgen.rendering.components.rating_badge import render_rating_badge
from reportgen.rendering.geometry import Box
from reportgen.rendering.theme import BrandTheme
from reportgen.schemas.blocks import MetricsBlock
from reportgen.schemas.report import ReportSpec
from reportgen.schemas.slides import SlideSpec


def render_cover_slide(
    slide: Any,
    slide_spec: SlideSpec,
    report_spec: ReportSpec,
    theme: BrandTheme,
    runtime: Any,
) -> None:
    apply_background(slide, theme, runtime)
    canvas_w = theme.canvas.width_in
    canvas_h = theme.canvas.height_in
    meta = report_spec.metadata

    # ════════════════════════════════════════════════════════════════
    # SECTION 1: NAVY HEADER BLOCK (top ~35% of slide)
    # Matches the HTML .report-header
    # ════════════════════════════════════════════════════════════════
    header_h = canvas_h * 0.34
    _add_filled_rect(slide, runtime, 0, 0, canvas_w, header_h, theme.palette.primary)

    # Orange accent triangle/wedge on the right (decorative)
    try:
        from pptx.util import Inches, Pt, Emu
        from pptx.oxml.ns import qn
        import lxml.etree as etree

        # Add a subtle gradient overlay effect on the right side
        wedge_w = 2.5
        wedge = slide.shapes.add_shape(
            runtime.MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            runtime.Inches(canvas_w - wedge_w),
            runtime.Inches(0),
            runtime.Inches(wedge_w),
            runtime.Inches(header_h),
        )
        wedge.fill.solid()
        wedge.fill.fore_color.rgb = _hex_to_rgb(runtime, theme.palette.accent)
        wedge.line.fill.background()
        # Set transparency to make it subtle
        wedge.fill.fore_color.brightness = 0.0
    except Exception:
        pass

    # ── Row 1: Logo + Firm Brand (left) | Report Type Badge (right) ──
    logo_path = Path(theme.logo_path) if theme.logo_path else None
    brand_left = 0.5
    if logo_path and logo_path.exists():
        try:
            # White chip behind logo
            chip = slide.shapes.add_shape(
                runtime.MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
                runtime.Inches(0.4),
                runtime.Inches(0.25),
                runtime.Inches(2.6),
                runtime.Inches(0.65),
            )
            chip.fill.solid()
            chip.fill.fore_color.rgb = _hex_to_rgb(runtime, theme.palette.background)
            chip.line.fill.background()
            slide.shapes.add_picture(
                str(logo_path),
                runtime.Inches(0.5),
                runtime.Inches(0.3),
                height=runtime.Inches(0.50),
            )
            brand_left = 1.3
        except Exception:
            pass

    # Firm name
    _add_text_in_rect(
        slide, runtime, brand_left, 0.3, 3.0, 0.35,
        theme.firm_name,
        theme.title_font.family, 16, True,
        "#FFFFFF",
        runtime.PP_ALIGN.LEFT,
    )
    # Firm tagline
    _add_text_in_rect(
        slide, runtime, brand_left, 0.65, 3.0, 0.2,
        theme.firm_tagline.upper(),
        theme.body_font.family, 8, False,
        "#A0B0C8",  # muted white on navy
        runtime.PP_ALIGN.LEFT,
    )

    # Report type badge (right side)
    report_type = f"{meta.report_type or 'Initiation'}"
    badge_box = slide.shapes.add_shape(
        runtime.MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        runtime.Inches(canvas_w - 4.2),
        runtime.Inches(0.3),
        runtime.Inches(3.5),
        runtime.Inches(0.3),
    )
    badge_box.fill.solid()
    badge_box.fill.fore_color.rgb = runtime.RGBColor(255, 255, 255)
    badge_box.fill.fore_color.brightness = 0.85  # semi-transparent
    badge_box.line.fill.background()
    tf = badge_box.text_frame
    tf.margin_left = runtime.Inches(0.05)
    tf.margin_top = runtime.Inches(0.02)
    p = tf.paragraphs[0]
    p.alignment = runtime.PP_ALIGN.CENTER
    r = p.add_run()
    r.text = report_type.upper()
    r.font.name = theme.body_font.family
    r.font.size = runtime.Pt(9)
    r.font.bold = False
    r.font.color.rgb = _hex_to_rgb(runtime, "#FFFFFF")

    # ── Row 2: Company Name + Rating Badge + Sector Tag ──
    company_name = report_spec.company.name
    _add_text_in_rect(
        slide, runtime, 0.5, 1.0, canvas_w * 0.6, 0.6,
        company_name,
        theme.title_font.family, 30, True,
        "#FFFFFF",
        runtime.PP_ALIGN.LEFT,
    )

    # Rating badge (BUY/HOLD/SELL)
    rating = (meta.rating or "").upper()
    if rating:
        rating_color = getattr(theme.rating_colors, rating, theme.palette.accent)
        render_rating_badge(
            slide,
            Box(left=0.5 + len(company_name) * 0.18, top=1.15, width=1.0, height=0.35),
            rating,
            theme,
            runtime,
        )

    # Sector/Exchange tag
    exchange_info = f"{report_spec.company.exchange}: {report_spec.company.ticker}"
    if report_spec.company.sector:
        exchange_info += f"  |  {report_spec.company.sector}"
    _add_text_in_rect(
        slide, runtime, 0.5, 1.6, canvas_w - 1.0, 0.2,
        exchange_info,
        theme.body_font.family, 9, False,
        "#8899AA",
        runtime.PP_ALIGN.LEFT,
    )

    # ── Row 3: Stats bar (CMP | TP | Market Cap | SAARTHI | Upside) ──
    stats_top = 1.95
    stats_height = 0.65

    # Divider line above stats
    _add_filled_rect(slide, runtime, 0.5, stats_top - 0.02, canvas_w - 1.0, 0.015, "#4A6AA0")

    stats_data = [
        ("CMP", _format_currency_compact(meta.cmp), f"As of {_format_date_ddmmyyyy(str(meta.report_date))}"),
        ("TARGET PRICE", _format_currency_compact(meta.target_price), ""),
    ]
    if meta.market_cap:
        stats_data.append(("MARKET CAP", _format_market_cap(meta.market_cap), ""))

    # Calculate upside
    try:
        cmp_val = float(meta.cmp)
        tp_val = float(meta.target_price)
        if cmp_val > 0:
            upside = ((tp_val - cmp_val) / cmp_val) * 100
            stats_data.append(("UPSIDE", f"+{upside:.0f}%", ""))
    except (ValueError, TypeError):
        pass

    stat_count = len(stats_data)
    stat_w = (canvas_w - 1.0) / stat_count

    for idx, (label, value, sub) in enumerate(stats_data):
        x = 0.5 + idx * stat_w

        # Label
        _add_text_in_rect(
            slide, runtime, x, stats_top, stat_w - 0.2, 0.15,
            label,
            theme.body_font.family, 7, False,
            "#8899AA",
            runtime.PP_ALIGN.LEFT,
        )
        # Value
        is_accent = idx == 0  # CMP in accent color
        _add_text_in_rect(
            slide, runtime, x, stats_top + 0.15, stat_w - 0.2, 0.3,
            value,
            theme.body_font.family, 16, True,
            theme.palette.accent if is_accent else "#FFFFFF",
            runtime.PP_ALIGN.LEFT,
        )
        # Sub-text
        if sub:
            _add_text_in_rect(
                slide, runtime, x, stats_top + 0.45, stat_w - 0.2, 0.15,
                sub,
                theme.body_font.family, 7, False,
                "#7788AA",
                runtime.PP_ALIGN.LEFT,
            )

        # Vertical separator
        if idx < stat_count - 1:
            _add_filled_rect(slide, runtime, x + stat_w - 0.15, stats_top + 0.05, 0.01, 0.5, "#3A5BA0")

    # Report date (bottom-right of header)
    _add_text_in_rect(
        slide, runtime, canvas_w - 3.5, header_h - 0.25, 3.0, 0.2,
        f"Report: {_format_date_ddmmyyyy(str(meta.report_date))}",
        theme.body_font.family, 8, False,
        "#8899AA",
        runtime.PP_ALIGN.RIGHT,
    )

    # ════════════════════════════════════════════════════════════════
    # SECTION 2: ORANGE TAGLINE BAND
    # Matches the HTML .tagline-band
    # ════════════════════════════════════════════════════════════════
    tagline_top = header_h
    tagline_h = 0.4
    _add_filled_rect(slide, runtime, 0, tagline_top, canvas_w, tagline_h, theme.palette.accent)

    # Get tagline from subtitle or a thesis summary
    tagline_text = slide_spec.subtitle or ""
    if not tagline_text:
        # Fallback: find text block in cover for tagline
        from reportgen.schemas.blocks import TextBlock
        for block in slide_spec.blocks:
            if isinstance(block, TextBlock):
                tagline_text = block.content[:200]
                break
    if tagline_text:
        _add_text_in_rect(
            slide, runtime, 0.5, tagline_top, canvas_w - 1.0, tagline_h,
            f'"{tagline_text}"' if not tagline_text.startswith('"') else tagline_text,
            theme.title_font.family, 12, True,
            theme.palette.primary,
            runtime.PP_ALIGN.LEFT,
        )

    # ════════════════════════════════════════════════════════════════
    # SECTION 3: COVER METRIC CARDS (bottom section)
    # Matches the HTML cover metric cards
    # ════════════════════════════════════════════════════════════════
    content_top = tagline_top + tagline_h + 0.1
    _render_cover_metrics(slide, slide_spec, theme, runtime, canvas_h, content_top)

    # ════════════════════════════════════════════════════════════════
    # SECTION 4: FOOTER
    # ════════════════════════════════════════════════════════════════
    footer_h = 0.35
    footer_top = canvas_h - footer_h
    _add_filled_rect(slide, runtime, 0, footer_top, canvas_w, footer_h, theme.palette.primary)

    # Footer left: SEBI
    _add_text_in_rect(
        slide, runtime, 0.5, footer_top, 3.5, footer_h,
        theme.footer_line,
        theme.footer_font.family, theme.footer_font.size_pt, False,
        theme.footer_font.color_hex,
        runtime.PP_ALIGN.LEFT,
    )

    # Footer center: Company name
    _add_text_in_rect(
        slide, runtime, canvas_w * 0.3, footer_top, canvas_w * 0.4, footer_h,
        f"{company_name} — Equity Research",
        theme.footer_font.family, theme.footer_font.size_pt, True,
        theme.palette.accent,
        runtime.PP_ALIGN.CENTER,
    )

    # Footer right: Date + Page 1
    formatted_date = _format_date_ddmmyyyy(str(meta.report_date))
    total_slides = 1  # Cover is page 1
    _add_text_in_rect(
        slide, runtime, canvas_w - 2.5, footer_top, 2.0, footer_h,
        f"{formatted_date}  |  Page 1",
        theme.footer_font.family, theme.footer_font.size_pt, False,
        theme.footer_font.color_hex,
        runtime.PP_ALIGN.RIGHT,
    )


def _render_cover_metrics(
    slide: Any,
    slide_spec: SlideSpec,
    theme: BrandTheme,
    runtime: Any,
    canvas_h: float,
    content_top: float,
) -> None:
    metrics_block = next(
        (b for b in slide_spec.blocks if isinstance(b, MetricsBlock)),
        None,
    )
    if not metrics_block or not metrics_block.items:
        return

    items = list(metrics_block.items)
    n = len(items)
    canvas_w = theme.canvas.width_in
    side = 0.5
    gap = 0.2
    total_w = canvas_w - 2 * side
    card_w = (total_w - gap * (n - 1)) / n
    card_h = 1.0
    card_top = content_top + 0.1

    # Color rotation for metric cards (matching HTML .metric-card variants)
    bg_colors = [
        theme.palette.primary,     # navy
        theme.palette.accent,      # orange
        theme.palette.secondary,   # mid-blue
        theme.palette.teal,        # teal
        theme.palette.primary,     # navy
        theme.palette.accent,      # orange
    ]

    for index, item in enumerate(items):
        left = side + index * (card_w + gap)
        bg = bg_colors[index % len(bg_colors)]

        card = slide.shapes.add_shape(
            runtime.MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            runtime.Inches(left),
            runtime.Inches(card_top),
            runtime.Inches(card_w),
            runtime.Inches(card_h),
        )
        card.fill.solid()
        card.fill.fore_color.rgb = _hex_to_rgb(runtime, bg)
        card.line.fill.background()

        # Determine text colors based on background
        is_accent_bg = bg == theme.palette.accent
        label_color = theme.palette.primary if is_accent_bg else "#C0C8D4"
        value_color = theme.palette.primary if is_accent_bg else "#FFFFFF"

        tf = card.text_frame
        tf.margin_left = runtime.Inches(0.08)
        tf.margin_right = runtime.Inches(0.08)
        tf.margin_top = runtime.Inches(0.08)
        tf.margin_bottom = runtime.Inches(0.05)

        # Label
        p_label = tf.paragraphs[0]
        p_label.alignment = runtime.PP_ALIGN.CENTER
        rl = p_label.add_run()
        rl.text = item.label.upper()
        rl.font.name = theme.body_font.family
        rl.font.size = runtime.Pt(8)
        rl.font.bold = False
        rl.font.color.rgb = _hex_to_rgb(runtime, label_color)

        # Value
        p_val = tf.add_paragraph()
        p_val.alignment = runtime.PP_ALIGN.CENTER
        rv = p_val.add_run()
        rv.text = str(item.value) if item.value is not None else ""
        rv.font.name = theme.metric_font.family
        rv.font.size = runtime.Pt(18)
        rv.font.bold = True
        rv.font.color.rgb = _hex_to_rgb(runtime, value_color)
